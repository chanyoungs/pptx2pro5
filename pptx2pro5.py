import os
import ntpath
from uuid import uuid4
from base64 import b64encode
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tkinter import *
from tkinter import ttk, filedialog
from zipfile import ZipFile
import shutil


def hex_to_rgb(hex):
    return tuple(int(hex.lstrip('#')[i:i+2], 16)/255. for i in (0, 2, 4))


def make_uuid():
    return str(uuid4()).upper()


def standardConversion(text):
    text = text.replace(u'[]', u'')
    text = text.replace(u'\u2013', u'-')
    text = text.replace(u'\u2014', u'-')
    text = text.replace(u'\u2018', u"'")
    text = text.replace(u'\u2019', u"'")
    text = text.replace(u'\u201c', u'"')
    text = text.replace(u'\u201d', u'"')
    text = text.replace(u'\u2026', u'...')
    text = text.replace(u'\xa0', u' ')  # nbsp
    return text


def get_filename(path):
    head, tail = ntpath.split(path)
    with_extension = tail or ntpath.basename(head)
    return os.path.splitext(with_extension)[0]


class PPTX2PRO():

    def __init__(self, path_to_template_pptx):

        # Load Template pptx
        self.pptx_template = Presentation(path_to_template_pptx)

        # Get Template dimensions in pixels
        pptx = self.pptx_template
        self.width = pptx.slide_width.inches * 96
        self.height = pptx.slide_height.inches * 96

        # Get Template background color in (R, G, B)
        slide = pptx.slides[0]

        if slide.background.fill.fore_color.type == MSO_COLOR_TYPE.RGB:
            self.background_color = hex_to_rgb(
                str(slide.background.fill.fore_color.rgb))
        else:
            self.background_color = (0, 0, 0)

        # Get shape positions and dimensions in pixels
        shape = slide.shapes[0]
        self.shape = {
            "left": shape.left.inches * 96,
            "top": shape.top.inches * 96,
            "width": shape.width.inches * 96,
            "height": shape.height.inches * 96
        }

        # Get text vertical alignment, margins, word-wrap
        text_frame = shape.text_frame
        self.text_frame = {
            # ENUM https://python-pptx.readthedocs.io/en/latest/api/enum/MsoVerticalAnchor.html#msoverticalanchor
            "v_align": text_frame.vertical_anchor,
            "margin": {
                "left": text_frame.margin_left.inches * 96,  # in pixels
                "right": text_frame.margin_right.inches * 96,
                "top": text_frame.margin_top.inches * 96,
                "bottom": text_frame.margin_bottom.inches * 96
            },
            "word_wrap": text_frame.word_wrap  # Boolean
        }

        # Get horizontal alignment
        paragraph = text_frame.paragraphs[0]
        # ENUM https://python-pptx.readthedocs.io/en/latest/api/enum/PpParagraphAlignment.html#ppparagraphalignment
        self.h_align = paragraph.alignment,

        # Get font formats
        font = paragraph.runs[0].font

        if font.color.type == MSO_COLOR_TYPE.RGB:
            font_color = hex_to_rgb(font.color.rgb)
        else:
            font_color = (1, 1, 1)

        self.font = {
            "bold": font.bold,  # boolean
            "italic": font.italic,
            "underline": font.underline,
            "color": font_color,
            "name": font.name,
            "size": font.size.pt  # Point
        }

    def rtfdata_text(self, text):
        slide = "{\\rtf1\\ansi\\ansicpg1252\\cocoartf1347\\cocoasubrtf570" + "\n" + \
            "\\cocoascreenfonts1{\\fonttbl\\f0\\fnil\\fcharset0 " + self.font["name"] + ";}" + "\n" \
            "{\\colortbl;\\red" + str(self.font["color"][0]) + "\\green" + str(self.font["color"][1]) + "\\blue" + str(self.font["color"][2]) + ";}" + "\n" + \
            "\\pard\\tx560\\tx1120\\tx1680\\tx2240\\tx2800\\tx3360\\tx3920\\tx4480\\tx5040\\tx5600\\tx6160\\tx6720\\pardirnatural\\qc" + "\n" + \
            "\\f0\\fs" + str(self.font["size"] * 2) + " \\cf1  " + text + "}"
        return b64encode(slide.encode()).decode()

    def headers(self):
        return f'''<RVPresentationDocument height="{self.height}" width="{self.width}" versionNumber="500" docType="0" creatorCode="1349676880" lastDateUsed="2015-08-08T22:38:35" usedCount="0" category="Speaker" resourcesDirectory="" backgroundColor="{self.background_color[0]} {self.background_color[1]} {self.background_color[2]} 1" drawingBackgroundColor="1" notes="" artist="" author="" album="" CCLIDisplay="0" CCLIArtistCredits="" CCLISongTitle="" CCLIPublisher="" CCLICopyrightInfo="" CCLILicenseNumber="" chordChartPath="">
	<timeline timeOffSet="0" selectedMediaTrackIndex="0" unitOfMeasure="60" duration="0" loop="0">
		<timeCues containerClass="NSMutableArray" />
		<mediaTracks containerClass="NSMutableArray" />
	</timeline>
	<bibleReference containerClass="NSMutableDictionary" />
	<_-RVProTransitionObject-_transitionObject transitionType="-1" transitionDuration="1" motionEnabled="0" motionDuration="20" motionSpeed="100" />
	<groups containerClass="NSMutableArray">
		<RVSlideGrouping name="" uuid="3AFCBE29-AC33-496E-A181-E7C4B4618FCB" color="0 0 0 0" serialization-array-index="0">
			<slides containerClass="NSMutableArray">'''

    def slide_text(self, index, name, text):
        return f'''
                    <RVDisplaySlide backgroundColor="{self.background_color[0]} {self.background_color[1]} {self.background_color[2]} 1" enabled="1" highlightColor="0 0 0 0" hotKey="" label="{name}" notes="" slideType="1" sort_index="' + sIndex + '" UUID="{make_uuid()}" drawingBackgroundColor="0" chordChartPath="" serialization-array-index="{index}">
						<cues containerClass="NSMutableArray" />
						<displayElements containerClass="NSMutableArray">
							<RVTextElement displayDelay="0" displayName="" locked="0" persistent="0" typeID="0" fromTemplate="0" bezelRadius="0" drawingFill="0" drawingShadow="1" drawingStroke="0" fillColor="0 0 0 0" rotation="0" source="" adjustsHeightToFit="0" verticalAlignment="0" RTFData="{self.rtfdata_text(text)}" revealType="0" serialization-array-index="0">
								<_-RVRect3D-_position x="0" y="0" z="0" width="{self.width}" height="{self.height}" />
								<_-D-_serializedShadow containerClass="NSMutableDictionary">
									<NSMutableString serialization-native-value="{{5, -5}}" serialization-dictionary-key="shadowOffset" />
									<NSNumber serialization-native-value="0" serialization-dictionary-key="shadowBlurRadius" />
									<NSColor serialization-native-value="0 0 0 0.3333333432674408" serialization-dictionary-key="shadowColor" />
								</_-D-_serializedShadow>
								<stroke containerClass="NSMutableDictionary">
									<NSColor serialization-native-value="0 0 0 1" serialization-dictionary-key="RVShapeElementStrokeColorKey" />
									<NSNumber serialization-native-value="1" serialization-dictionary-key="RVShapeElementStrokeWidthKey" />
								</stroke>
							</RVTextElement>
						</displayElements>
						<_-RVProTransitionObject-_transitionObject transitionType="-1" transitionDuration="1" motionEnabled="0" motionDuration="20" motionSpeed="100" />
					</RVDisplaySlide>'''

    def slide_image(self, index, filename):
        return f'''
                    <RVDisplaySlide backgroundColor="{self.background_color[0]} {self.background_color[1]} {self.background_color[2]} 1" highlightColor="" drawingBackgroundColor="0" enabled="1" hotKey="" label="" notes="" UUID="{make_uuid()}" chordChartPath="" serialization-array-index="{index}">
                        <cues containerClass="NSMutableArray">
                            <RVMediaCue UUID="DF44057F-EF40-48D7-B13F-0E8D7BE8C852" displayName="{filename}" enabled="1" timeStamp="0" delayTime="0" behavior="1" alignment="4" serialization-array-index="0" elementClassName="RVImageElement">
                                <element displayName="ImageSample1.jpg" displayDelay="0" locked="0" persistent="0" typeID="0" fromTemplate="0" bezelRadius="0" drawingFill="0" drawingShadow="0" drawingStroke="0" fillColor="1 1 1 1" rotation="0" source="{filename}" flippedHorizontally="0" flippedVertically="0" scaleBehavior="3" manufactureURL="" manufactureName="" format="">
                                    <_-RVRect3D-_position x="0" y="0" z="0" width="0" height="0" />
                                    <_-D-_serializedShadow containerClass="NSMutableDictionary">
                                        <NSNumber serialization-native-value="0" serialization-dictionary-key="shadowBlurRadius" />
                                        <NSColor serialization-native-value="0 0 0 0" serialization-dictionary-key="shadowColor" />
                                        <NSMutableString serialization-native-value="{{0, 0}}" serialization-dictionary-key="shadowOffset" />
                                    </_-D-_serializedShadow>
                                    <stroke containerClass="NSMutableDictionary">
                                        <NSColor serialization-dictionary-key="RVShapeElementStrokeColorKey" serialization-native-value="0 0 0 1" />
                                        <NSNumber serialization-dictionary-key="RVShapeElementStrokeWidthKey" serialization-native-value="1.0" />
                                    </stroke>
                                    <effects containerClass="NSMutableArray" />
                                </element>
                            </RVMediaCue>
                        </cues>
                        <displayElements containerClass="NSMutableArray" />
                    </RVDisplaySlide>'''

    def closure(self):
        return '''
				</slides>
			</RVSlideGrouping>
		</groups>
	<arrangements containerClass="NSMutableArray">
		<RVSongArrangement name="New Arrangement" uuid="8DD67BB3-30A2-4859-92CC-6B34181ADE5F" color="0 0 0 0" serialization-array-index="0">
			<groupIDs containerClass="NSMutableArray">
				<NSMutableString serialization-native-value="3AFCBE29-AC33-496E-A181-E7C4B4618FCB" serialization-array-index="0" />
			</groupIDs>
		</RVSongArrangement>
	</arrangements>
</RVPresentationDocument>'''

    def convert(self, text_mode, path_to_pptx, save_path):
        slides = Presentation(path_to_pptx).slides
        with open(os.path.join(save_path, get_filename(path_to_pptx) + ".pro5"), 'w', encoding="utf-8") as f:
            f.write(self.headers())

            if text_mode == True:
                for i in range(len(slides)):
                    t = ""
                    if len(slides[i].shapes) > 0:
                        t = slides[i].shapes[0].text

                    text = standardConversion(t.replace("\n", "\\\n"))
                    f.write(self.slide_text(index=i, name="", text=text))

            else:
                media_path = os.path.join("temp", "media")
                if not os.path.exists(media_path):
                    os.makedirs(media_path)
                for i in range(len(slides)):
                    shapes = slides[i].shapes
                    if len(shapes) > 0:
                        for shape in shapes:
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image = shape.image
                                image_file_name = f"{get_filename(path_to_pptx)}_{i+1}.{image.ext}"
                                f.write(self.slide_image(
                                    index=i, filename=image_file_name))
                                with open(os.path.join(media_path, image_file_name), 'wb') as img_file:
                                    img_file.write(image.blob)
                                break

            f.write(self.closure())


class Application(Frame):
    def __init__(self, master=None):
        super().__init__(master)
        with open("data/config.txt", "r") as f:
            self.location = f.read()
        self.master = master
        self.master.title("PowerPoint to ProPresenter Converter")
        self.master.geometry("1080x720")
        # self.master.maxsize(1000, 400)
        self.pack()
        self.create_widgets()
        self.files = None

    def create_widgets(self):
        if self.location == "":
            label_text = "Save location not yet set"
            label_color = "red"
            button_text = "Set save location"
        else:
            label_text = f"Save location set to '{self.location}'"
            label_color = "green"
            button_text = "Change save location"

        row = 0
        self.save_location_button = Button(
            self, text=button_text, font=("arial", 17), command=self.set_save_location, fg="white", bg="grey")
        self.save_location_button.grid(
            row=row, column=0, columnspan=2, pady=(50, 0))
        row += 1

        self.location_label = Label(self, font=(
            "arial", 12), text=label_text, fg=label_color)
        self.location_label.grid(row=row, column=0, columnspan=2, pady=(0, 50))
        row += 1

        self.open_template_button = Button(
            self, text="Open and edit template presentation", font=("arial", 17), command=self.open_template, fg="white", bg="grey")
        self.open_template_button.grid(
            row=row, column=0, columnspan=2, pady=(0, 50))
        row += 1

        self.open_button = Button(
            self, text="Select PowerPoint Presentation(s)", font=("arial", 17), command=self.open_powerpoint_files, fg="white", bg="grey")
        self.open_button.grid(row=row, column=0, columnspan=2)
        row += 1

        self.logs = StringVar()
        self.logs.set("No files selected")
        self.files_label = Label(
            self, textvariable=self.logs, fg="red", font=("arial", 12), justify=LEFT)
        self.files_label.grid(row=row, column=0, columnspan=2, pady=(0, 50))
        row += 1

        self.text_mode_boolean = BooleanVar()
        self.text_mode_boolean.set(True)
        self.text_mode_radiobutton = Radiobutton(
            self, text="Text Slides", variable=self.text_mode_boolean, value=True, command=self.slide_mode_command)
        self.text_mode_radiobutton.grid(row=row, column=0, pady=(0, 50))

        self.image_mode_radiobutton = Radiobutton(
            self, text="Image Slides", variable=self.text_mode_boolean, value=False, command=self.slide_mode_command)
        self.image_mode_radiobutton.grid(row=row, column=1, pady=(0, 50))
        row += 1

        self.pro5_boolean = BooleanVar()
        self.pro5_boolean.set(True)
        self.pro5_checkbutton = Checkbutton(
            self, text='.pro5', font=("arial", 12), variable=self.pro5_boolean, onvalue=True, offvalue=False)
        self.pro5_checkbutton.grid(row=row, column=0, pady=(0, 50))

        self.proBundle_boolean = BooleanVar()
        self.proBundle_boolean.set(False)
        self.proBundle_checkbutton = Checkbutton(
            self, text='.proBundle', font=("arial", 12), variable=self.proBundle_boolean, onvalue=True, offvalue=False)
        self.proBundle_checkbutton.grid(row=row, column=1, pady=(0, 50))

        self.checkbutton_row = row

        row += 1

        self.convert_button = Button(
            self, text="Convert!", font=("arial", 17), command=self.convert, fg="white", bg="grey")
        self.convert_button.grid(row=row, column=0, columnspan=2, pady=(0, 50))
        row += 1

        self.progress = ttk.Progressbar(
            self, orient=HORIZONTAL, mode='determinate', length=1000)
        self.progress.grid(row=row, column=0, columnspan=2)
        row += 1

        self.footer_text = StringVar()
        self.footer_label = Label(
            self, textvariable=self.footer_text, font=("arial", 12))
        self.footer_label.grid(row=row, column=0, columnspan=2, pady=(0, 50))
        row += 1

    def slide_mode_command(self):
        if self.text_mode_boolean.get() == True:
            self.pro5_checkbutton.grid(
                row=self.checkbutton_row, column=0, pady=(0, 50))
            self.proBundle_checkbutton.grid(
                row=self.checkbutton_row, column=1, pady=(0, 50))
        else:
            self.pro5_checkbutton.grid_forget()
            self.proBundle_checkbutton.grid_forget()
            self.footer_text.set("")

    def open_template(self):
        os.system("start " + "./data/Template.pptx")

    def open_powerpoint_files(self):
        self.files = filedialog.askopenfilenames(
            initialdir="/", title="Select PowerPoint Presentation(s)", filetypes=[("PowerPoint Presentations", "*.pptx")])

        if self.files != "":
            logs = f"{len(self.files)} files selected:"
            index = 1
            for i in range(len(self.files)):
                if (i > 4):
                    logs += "\n..."
                    break
                logs += f"\n\t{index}. {get_filename(self.files[i])}"
                index += 1
            self.logs.set(logs)
            self.files_label.config(fg="green")
            self.progress["value"] = 0
            self.progress["maximum"] = len(self.files)
            self.footer_text.set("")

    def set_save_location(self):
        location = filedialog.askdirectory(
            initialdir=self.location, title="Select save location")

        if location != "":
            with open("data/config.txt", "w", encoding="utf-8") as f:
                f.write(location)
            self.location = location
            self.location_label.config(
                text=f"Save location set to '{self.location}'", fg="green")
            self.save_location_button.config(text="Change save location")

    def convert(self):
        if self.location == "":
            self.footer_text.set("Set save location first!")
            self.footer_label.config(fg="red")
            return
        elif self.files is None:
            self.footer_text.set("Choose files first!")
            self.footer_label.config(fg="red")
            return
        elif self.pro5_boolean.get() == False and self.proBundle_boolean.get() == False and self.text_mode_boolean.get() == True:
            self.footer_text.set("Check at least one of .pro5 or .proBundle")
            self.footer_label.config(fg="red")
        else:
            self.footer_text.set("Converting...")
            self.footer_label.config(fg="black")
            if self.proBundle_boolean.get() == True or self.text_mode_boolean.get() == False:
                zip = ZipFile(os.path.join(
                    self.location, 'Slides.proBundle'), 'w')

            if self.pro5_boolean.get() == True and self.text_mode_boolean.get() == True:
                save_path = self.location
            else:
                temp_path = "./temp"
                save_path = temp_path
                if not os.path.exists(save_path):
                    os.makedirs(save_path)

            for i in range(len(self.files)):
                pptx2pro = PPTX2PRO(
                    path_to_template_pptx="./data/Template.pptx")
                pptx2pro.convert(text_mode=self.text_mode_boolean.get(),
                                 path_to_pptx=self.files[i], save_path=save_path)
                self.progress["value"] = i + 1
                self.progress.update()
                if self.proBundle_boolean.get() == True or self.text_mode_boolean.get() == False:
                    zip.write(os.path.join(save_path, get_filename(
                        self.files[i]) + ".pro5"), os.path.join(get_filename(self.files[i]) + ".pro5"))
            if self.text_mode_boolean.get() == False:
                images = os.listdir(os.path.join("temp", "media"))
                for image in images:
                    zip.write(os.path.join(save_path, "media", image),
                              os.path.join("media", image))

            self.footer_text.set("COMPLETE!")
            self.footer_label.config(fg="green")

            if self.proBundle_boolean.get() == True or self.text_mode_boolean.get() == False:
                zip.close()
                shutil.rmtree(temp_path)


root = Tk()
app = Application(master=root)

app.mainloop()
